"""
RAG Service using LangChain and Gemini.
Supports PDF, PPTX, DOCX, TXT, and code files.
"""

import os
import io
import tempfile
from typing import List, Dict, Any, Optional
from dataclasses import dataclass

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_core.documents import Document
from supabase import Client


@dataclass
class ChunkResult:
    """Search result."""
    id: str
    material_id: str
    chunk_text: str
    chunk_index: int
    file_name: str
    page_number: Optional[int]
    category: Optional[str]
    topic: Optional[str]
    week_number: Optional[int]
    similarity: float


@dataclass
class RAGResponse:
    """RAG answer response."""
    answer: str
    sources: List[Dict[str, Any]]


class RAGService:
    """RAG Service for intelligent search with multi-format support."""
    
    # Supported file types
    SUPPORTED_TYPES = {'pdf', 'pptx', 'docx', 'doc', 'txt', 'md', 'py', 'js', 'ts', 'cpp', 'c', 'java', 'html', 'css', 'json', 'yaml', 'yml'}
    
    def __init__(self, supabase_client: Client, gemini_api_key: str):
        self.supabase = supabase_client
        self.gemini_api_key = gemini_api_key
        
        self.embeddings = GoogleGenerativeAIEmbeddings(
            model="models/text-embedding-004",
            google_api_key=gemini_api_key,
            task_type="retrieval_document"
        )
        
        self.llm = ChatGoogleGenerativeAI(
            model="gemini-2.5-flash",
            google_api_key=gemini_api_key,
            temperature=0.3,
        )
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
    
    def extract_text(self, file_content: bytes, file_type: str, file_name: str) -> List[Document]:
        """
        Extract text from various file formats.
        Returns list of Document objects with page/slide metadata.
        """
        file_type = file_type.lower()
        
        if file_type == 'pdf':
            return self._extract_pdf(file_content)
        elif file_type == 'pptx':
            return self._extract_pptx(file_content, file_name)
        elif file_type in ('docx', 'doc'):
            return self._extract_docx(file_content, file_name)
        elif file_type in ('txt', 'md', 'py', 'js', 'ts', 'cpp', 'c', 'java', 'html', 'css', 'json', 'yaml', 'yml'):
            return self._extract_text(file_content, file_name, file_type)
        else:
            # Try to read as plain text
            try:
                return self._extract_text(file_content, file_name, file_type)
            except:
                return []
    
    def _extract_pdf(self, file_content: bytes) -> List[Document]:
        """Extract text from PDF."""
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name
        
        try:
            loader = PyPDFLoader(tmp_path)
            return loader.load()
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    
    def _extract_pptx(self, file_content: bytes, file_name: str) -> List[Document]:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(file_content))
            documents = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                
                if texts:
                    content = "\n".join(texts)
                    documents.append(Document(
                        page_content=content,
                        metadata={"page": slide_num, "source": file_name, "type": "slide"}
                    ))
            
            return documents
        except Exception as e:
            print(f"PPTX extraction error: {e}")
            return []
    
    def _extract_docx(self, file_content: bytes, file_name: str) -> List[Document]:
        """Extract text from Word document."""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(io.BytesIO(file_content))
            paragraphs = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paragraphs.append(cell.text.strip())
            
            if paragraphs:
                content = "\n\n".join(paragraphs)
                return [Document(
                    page_content=content,
                    metadata={"page": 1, "source": file_name, "type": "document"}
                )]
            return []
        except Exception as e:
            print(f"DOCX extraction error: {e}")
            return []
    
    def _extract_text(self, file_content: bytes, file_name: str, file_type: str) -> List[Document]:
        """Extract text from plain text/code files."""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    content = file_content.decode(encoding)
                    break
                except:
                    continue
            else:
                content = file_content.decode('utf-8', errors='ignore')
            
            if content.strip():
                return [Document(
                    page_content=content,
                    metadata={"page": 1, "source": file_name, "type": file_type}
                )]
            return []
        except Exception as e:
            print(f"Text extraction error: {e}")
            return []
    
    def _enrich_with_metadata(self, chunk_text: str, metadata: Dict[str, Any]) -> str:
        """
        Enrich chunk text with course material metadata for better search.
        This helps the search find content by title, topic, description etc.
        """
        enrichment_parts = []
        
        if metadata.get("title"):
            enrichment_parts.append(f"Title: {metadata['title']}")
        if metadata.get("description"):
            enrichment_parts.append(f"Description: {metadata['description']}")
        if metadata.get("topic"):
            enrichment_parts.append(f"Topic: {metadata['topic']}")
        if metadata.get("category"):
            enrichment_parts.append(f"Category: {metadata['category']}")
        if metadata.get("tags"):
            tags = metadata['tags']
            if isinstance(tags, list):
                enrichment_parts.append(f"Tags: {', '.join(tags)}")
        if metadata.get("week_number"):
            enrichment_parts.append(f"Week: {metadata['week_number']}")
        
        if enrichment_parts:
            enrichment = "[METADATA]\n" + "\n".join(enrichment_parts) + "\n[CONTENT]\n"
            return enrichment + chunk_text
        
        return chunk_text
    
    def chunk_documents(self, documents: List[Document]) -> List[Document]:
        """Split documents into chunks."""
        return self.text_splitter.split_documents(documents)
    
    def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Generate embeddings for texts."""
        return self.embeddings.embed_documents(texts)
    
    def generate_query_embedding(self, query: str) -> List[float]:
        """Generate embedding for query."""
        query_emb = GoogleGenerativeAIEmbeddings(
            model="models/text-embedding-004",
            google_api_key=self.gemini_api_key,
            task_type="retrieval_query"
        )
        return query_emb.embed_query(query)
    
    async def index_material(
        self, 
        material_id: str, 
        file_content: bytes,
        file_name: str,
        file_type: str = "pdf",
        title: Optional[str] = None,
        description: Optional[str] = None,
        category: Optional[str] = None,
        topic: Optional[str] = None,
        week_number: Optional[int] = None,
        tags: Optional[List[str]] = None
    ) -> Dict[str, Any]:
        """
        Index a file into vector database with metadata enrichment.
        Supports PDF, PPTX, DOCX, TXT, and code files.
        """
        file_type = file_type.lower()
        
        # Check if supported
        if file_type not in self.SUPPORTED_TYPES:
            return {"success": False, "error": f"File type '{file_type}' not supported for indexing", "chunks_created": 0}
        
        try:
            # Delete existing chunks
            self.supabase.table("document_chunks").delete().eq(
                "material_id", material_id
            ).execute()
            
            # Extract text based on file type
            documents = self.extract_text(file_content, file_type, file_name)
            if not documents:
                return {"success": False, "error": "No text extracted from file", "chunks_created": 0}
            
            # Chunk
            chunks = self.chunk_documents(documents)
            if not chunks:
                return {"success": False, "error": "No chunks created", "chunks_created": 0}
            
            # Prepare metadata for enrichment
            metadata = {
                "title": title,
                "description": description,
                "category": category,
                "topic": topic,
                "week_number": week_number,
                "tags": tags
            }
            
            # Enrich chunks with metadata and generate embeddings
            enriched_texts = []
            for chunk in chunks:
                enriched = self._enrich_with_metadata(chunk.page_content, metadata)
                enriched_texts.append(enriched)
            
            embeddings = self.generate_embeddings(enriched_texts)
            
            # Store in database
            records = []
            for idx, (chunk, emb, enriched_text) in enumerate(zip(chunks, embeddings, enriched_texts)):
                page = chunk.metadata.get("page")
                if page is not None:
                    page = int(page)
                    # PDF pages are 0-indexed, others might be 1-indexed
                    if file_type == "pdf":
                        page += 1
                
                records.append({
                    "material_id": material_id,
                    "chunk_text": enriched_text,  # Store enriched text for search
                    "chunk_index": idx,
                    "embedding": emb,
                    "file_name": file_name,
                    "page_number": page,
                    "category": category,
                    "topic": topic,
                    "week_number": week_number
                })
            
            self.supabase.table("document_chunks").insert(records).execute()
            
            # Mark as indexed
            self.supabase.table("course_materials").update({
                "is_indexed": True
            }).eq("id", material_id).execute()
            
            return {"success": True, "chunks_created": len(records), "message": f"Indexed {len(records)} chunks from {file_type.upper()}"}
            
        except Exception as e:
            return {"success": False, "error": str(e), "chunks_created": 0}
    
    async def search(
        self, 
        query: str, 
        limit: int = 5,
        threshold: float = 0.5,
        category: Optional[str] = None,
        week: Optional[int] = None
    ) -> List[ChunkResult]:
        """Vector similarity search."""
        embedding = self.generate_query_embedding(query)
        
        response = self.supabase.rpc(
            "match_documents",
            {
                "query_embedding": embedding,
                "match_threshold": threshold,
                "match_count": limit,
                "filter_category": category,
                "filter_week": week
            }
        ).execute()
        
        results = []
        for row in response.data:
            # Extract just the content part for display (remove metadata prefix)
            chunk_text = row["chunk_text"]
            if "[CONTENT]" in chunk_text:
                chunk_text = chunk_text.split("[CONTENT]", 1)[1].strip()
            
            results.append(ChunkResult(
                id=row["id"],
                material_id=row["material_id"],
                chunk_text=chunk_text,
                chunk_index=row["chunk_index"],
                file_name=row["file_name"],
                page_number=row["page_number"],
                category=row["category"],
                topic=row["topic"],
                week_number=row["week_number"],
                similarity=row["similarity"]
            ))
        
        return results
    
    async def ask(
        self, 
        question: str, 
        limit: int = 5,
        category: Optional[str] = None,
        week: Optional[int] = None
    ) -> RAGResponse:
        """Full RAG pipeline."""
        chunks = await self.search(question, limit, 0.4, category, week)
        
        if not chunks:
            return RAGResponse(
                answer="I couldn't find relevant information in the course materials.",
                sources=[]
            )
        
        # Build context
        context_parts = []
        for i, c in enumerate(chunks):
            src = f"[Source {i+1}: {c.file_name}"
            if c.page_number:
                src += f", Page {c.page_number}"
            if c.topic:
                src += f", Topic: {c.topic}"
            src += "]"
            context_parts.append(f"{src}\n{c.chunk_text}")
        
        context = "\n\n---\n\n".join(context_parts)
        
        # Generate answer
        prompt = f"""You are an AI tutor helping students. Answer using ONLY the provided context from course materials. Cite sources with [Source X].

CONTEXT:
{context}

QUESTION: {question}

ANSWER:"""

        response = self.llm.invoke(prompt)
        
        sources = [{
            "file_name": c.file_name,
            "page_number": c.page_number,
            "excerpt": c.chunk_text[:300] + "..." if len(c.chunk_text) > 300 else c.chunk_text,
            "similarity": round(c.similarity, 3),
            "material_id": c.material_id
        } for c in chunks]
        
        return RAGResponse(answer=response.content, sources=sources)
